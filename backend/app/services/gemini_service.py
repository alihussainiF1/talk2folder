"""
Gemini Service - Direct file upload and chat with Gemini API.

This provides a "fast path" for small files (<10MB) that bypasses the 
traditional RAG pipeline (chunk -> embed -> store -> retrieve).

Instead, files are uploaded directly to Gemini's File API and queried natively,
which is:
- Much faster for small files
- Better at understanding document structure
- Supports native PDF, DOCX, etc. understanding
"""

import asyncio
import json
from typing import AsyncIterator
from functools import partial
from concurrent.futures import ThreadPoolExecutor
import ssl
import http.client

import google.generativeai as genai
from tenacity import retry, stop_after_attempt, wait_exponential_jitter, retry_if_exception_type
from google.generativeai.types import HarmCategory, HarmBlockThreshold

from app.config import get_settings

settings = get_settings()

# Configure Gemini
genai.configure(api_key=settings.google_api_key)

_executor = ThreadPoolExecutor(max_workers=4)

# Size thresholds for fast path vs RAG
MAX_FILE_SIZE_BYTES = 10 * 1024 * 1024  # 10MB per file
MAX_TOTAL_SIZE_BYTES = 100 * 1024 * 1024  # 100MB total for fast path
MAX_FILES_FOR_FAST_PATH = 50  # Max files before switching to RAG


# MIME types that Gemini supports natively
GEMINI_NATIVE_MIMES = {
    "application/pdf",
    "text/plain",
    "text/markdown", 
    "text/html",
    "text/csv",
    "application/json",
    "text/x-python",
    "application/x-python-code",
    "text/javascript",
    "application/javascript",
    "image/png",
    "image/jpeg",
    "image/jpg",
    "image/gif",
    "image/webp",
}

OFFICE_MIMES_NEED_CONVERSION = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _extract_text_from_docx(content: bytes) -> str:
    """Extract text from DOCX file."""
    import io
    from docx import Document
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_text_from_xlsx(content: bytes) -> str:
    """Extract text from XLSX file."""
    import io
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PPTX file."""
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n\n".join(lines)


def convert_office_to_text(content: bytes, mime_type: str, filename: str) -> tuple[bytes, str, str]:
    """
    Convert Office files to plain text for Gemini compatibility.
    Returns (content, mime_type, filename) tuple.
    """
    try:
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = _extract_text_from_docx(content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            text = _extract_text_from_xlsx(content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = _extract_text_from_pptx(content)
        else:
            return content, mime_type, filename
        
        base_name = filename.rsplit(".", 1)[0] if "." in filename else filename
        return text.encode("utf-8"), "text/plain", f"{base_name}.txt"
    except Exception as e:
        print(f"[Gemini] Failed to convert {filename}: {e}, uploading as-is")
        return content, mime_type, filename


def _get_file_size(f: dict) -> int:
    """Safely get file size as integer (Drive API returns strings)."""
    size = f.get("size", 0)
    if size is None:
        return 0
    try:
        return int(size)
    except (ValueError, TypeError):
        return 0


def should_use_fast_path(files: list[dict]) -> bool:
    """
    Determine if we should use the fast Gemini File API path.
    
    Returns True if:
    - Total number of files <= MAX_FILES_FOR_FAST_PATH
    - Total size <= MAX_TOTAL_SIZE_BYTES
    - All files are supported by Gemini natively
    """
    if len(files) > MAX_FILES_FOR_FAST_PATH:
        return False
    
    total_size = sum(_get_file_size(f) for f in files)
    if total_size > MAX_TOTAL_SIZE_BYTES:
        return False
    
    # Check if any file is too large
    for f in files:
        if _get_file_size(f) > MAX_FILE_SIZE_BYTES:
            return False
    
    return True


@retry(
    stop=stop_after_attempt(5),
    wait=wait_exponential_jitter(initial=2, max=30, jitter=5),
    reraise=True,
    retry=retry_if_exception_type((ssl.SSLError, http.client.RemoteDisconnected, ConnectionError, TimeoutError, OSError)),
    before_sleep=lambda retry_state: print(f"[Gemini] Retry {retry_state.attempt_number}/5 for upload after {type(retry_state.outcome.exception()).__name__}"),
)
def _sync_upload_file(file_content: bytes, filename: str, mime_type: str) -> str:
    """Synchronously upload a file to Gemini File API. Returns file URI."""
    import tempfile
    import os
    import time
    
    # Map Google Workspace types to export types for filename
    extension_map = {
        "application/vnd.google-apps.document": ".docx",
        "application/vnd.google-apps.spreadsheet": ".xlsx",
        "application/vnd.google-apps.presentation": ".pptx",
        "application/pdf": ".pdf",
        "text/plain": ".txt",
        "text/markdown": ".md",
        "text/csv": ".csv",
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/jpg": ".jpg",
        "image/gif": ".gif",
        "image/webp": ".webp",
        "image/heic": ".heic",
        "video/mp4": ".mp4",
        "video/quicktime": ".mov",
        "audio/mpeg": ".mp3",
        "audio/wav": ".wav",
    }
    
    # Ensure filename has extension
    if "." not in filename:
        ext = extension_map.get(mime_type, ".bin")
        filename = f"{filename}{ext}"
    
    # Write to temp file and upload with retry
    with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{filename}") as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name
    
    try:
        uploaded = genai.upload_file(tmp_path, mime_type=mime_type)
        return uploaded.uri
    except Exception as e:
        print(f"[Gemini] Upload failed: {e}")
        raise
    finally:
        os.unlink(tmp_path)


async def upload_file_to_gemini(file_content: bytes, filename: str, mime_type: str) -> str:
    """
    Upload a file to Gemini File API.
    
    Args:
        file_content: Raw file bytes
        filename: Name of the file
        mime_type: MIME type of the file
        
    Returns:
        Gemini file URI (e.g., "https://generativelanguage.googleapis.com/...")
    """
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(
        _executor,
        _sync_upload_file,
        file_content,
        filename,
        mime_type
    )


async def upload_files_to_gemini(files_with_content: list[tuple[dict, bytes]], max_concurrent: int = 1) -> list[dict]:
    """
    Upload multiple files to Gemini sequentially to avoid SSL errors.
    
    Args:
        files_with_content: List of (file_metadata, content_bytes) tuples
        max_concurrent: Maximum concurrent uploads (default 1 = sequential)
        
    Returns:
        List of file info dicts with gemini_uri added
    """
    semaphore = asyncio.Semaphore(max_concurrent)
    total = len([f for f, c in files_with_content if c is not None])
    completed = 0
    success_count = 0
    fail_count = 0
    
    async def upload_with_semaphore(file_meta: dict, content: bytes) -> dict:
        nonlocal completed, success_count, fail_count
        
        mime_type = file_meta.get("mimeType", "application/octet-stream")
        filename = file_meta["name"]
        
        if mime_type == "application/vnd.google-apps.document":
            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif mime_type == "application/vnd.google-apps.spreadsheet":
            mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif mime_type == "application/vnd.google-apps.presentation":
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        if mime_type in OFFICE_MIMES_NEED_CONVERSION:
            content, mime_type, filename = convert_office_to_text(content, mime_type, filename)
            print(f"[Gemini] Converted {file_meta['name']} -> {filename} (text)")
        
        async with semaphore:
            try:
                print(f"[Gemini] Uploading {filename}...")
                uri = await upload_file_to_gemini(content, filename, mime_type)
                completed += 1
                success_count += 1
                print(f"[Gemini] OK: {filename} ({completed}/{total})")
                return {
                    "id": file_meta["id"],
                    "name": file_meta["name"],
                    "path": file_meta.get("path", file_meta["name"]),
                    "mime_type": file_meta.get("mimeType", ""),
                    "gemini_uri": uri,
                }
            except Exception as e:
                completed += 1
                fail_count += 1
                print(f"[Gemini] FAILED: {file_meta['name']} ({completed}/{total}) - {e}")
                return {
                    "id": file_meta["id"],
                    "name": file_meta["name"],
                    "path": file_meta.get("path", file_meta["name"]),
                    "mime_type": file_meta.get("mimeType", ""),
                    "gemini_uri": None,
                }
    
    tasks = []
    for file_meta, content in files_with_content:
        if content is None:
            continue
        tasks.append(upload_with_semaphore(file_meta, content))
    
    print(f"[Gemini] Starting upload of {total} files (concurrency: {max_concurrent})")
    results = await asyncio.gather(*tasks)
    print(f"[Gemini] Upload complete: {success_count} succeeded, {fail_count} failed")
    
    return results


def _create_model():
    """Create a Gemini model instance with safety settings."""
    return genai.GenerativeModel(
        "gemini-3-flash-preview",
        safety_settings={
            HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
        },
        system_instruction="""You are a helpful assistant that answers questions about documents.
        
CRITICAL RULES:
1. Base your answers ONLY on the provided documents
2. ALWAYS cite your sources using the format [Source: filename]
3. If information is not in the documents, say "I couldn't find this in the provided documents"
4. Be concise but thorough
5. When quoting from documents, use quotation marks

When listing files, format them nicely and mention their types."""
    )


def _get_file_from_uri(uri: str):
    """Get a Gemini file object from its URI."""
    # URI format: https://generativelanguage.googleapis.com/v1beta/files/abc123
    # Extract file name from URI
    file_name = uri.split("/")[-1]
    return genai.get_file(file_name)


async def chat_with_files(
    message: str,
    gemini_files: list[dict],
    history: list[dict] | None = None,
) -> dict:
    """
    Chat with files using Gemini's native file understanding.
    
    Args:
        message: User's message
        gemini_files: List of file dicts with gemini_uri
        history: Previous conversation history
        
    Returns:
        Dict with content and citations
    """
    model = _create_model()
    
    # Build content parts with files
    parts = []
    file_names = []
    
    # Get actual file objects from Gemini
    for f in gemini_files:
        uri = f.get("gemini_uri")
        if uri:
            try:
                gemini_file = _get_file_from_uri(uri)
                parts.append(gemini_file)
                file_names.append(f["name"])
            except Exception as e:
                print(f"[Gemini] Failed to get file {f['name']}: {e}")
    
    # Add the user's question
    parts.append(f"Please analyze the document(s) above and answer: {message}")
    
    # Build chat history
    chat_history = []
    if history:
        for msg in history:
            role = "user" if msg["role"] == "user" else "model"
            chat_history.append({
                "role": role,
                "parts": [msg["content"]]
            })
    
    loop = asyncio.get_running_loop()
    
    def _sync_generate():
        if chat_history:
            chat = model.start_chat(history=chat_history)
            response = chat.send_message(parts)
        else:
            response = model.generate_content(parts)
        return response.text
    
    content = await loop.run_in_executor(_executor, _sync_generate)
    
    # Extract citations from the response - match files mentioned and include metadata
    citations = []
    for f in gemini_files:
        fname = f.get("name", "")
        if fname and fname.lower() in content.lower():
            citations.append({
                "file_name": fname,
                "file_id": f.get("id"),
                "drive_file_id": f.get("id"),
                "mime_type": f.get("mime_type", ""),
                "chunk_index": 0,
            })
    
    return {
        "content": content,
        "citations": citations if citations else None,
    }


async def stream_chat_with_files(
    message: str,
    gemini_files: list[dict],
    history: list[dict] | None = None,
) -> AsyncIterator[str]:
    """
    Stream chat responses with files using Gemini's native file understanding.
    
    Args:
        message: User's message
        gemini_files: List of file dicts with gemini_uri
        history: Previous conversation history
        
    Yields:
        Text chunks as they're generated
    """
    model = _create_model()
    
    # Build content parts with files
    parts = []
    file_names = []
    
    # Get actual file objects from Gemini
    for f in gemini_files:
        uri = f.get("gemini_uri")
        if uri:
            try:
                gemini_file = _get_file_from_uri(uri)
                parts.append(gemini_file)
                file_names.append(f["name"])
            except Exception as e:
                print(f"[Gemini] Failed to get file {f['name']}: {e}")
    
    # Add the user's question
    parts.append(f"Please analyze the document(s) above and answer: {message}")
    
    # Build chat history
    chat_history = []
    if history:
        for msg in history:
            role = "user" if msg["role"] == "user" else "model"
            chat_history.append({
                "role": role,
                "parts": [msg["content"]]
            })
    
    loop = asyncio.get_running_loop()
    
    import queue
    import threading
    
    q = queue.Queue()
    
    def _sync_stream():
        try:
            if chat_history:
                chat = model.start_chat(history=chat_history)
                response = chat.send_message(parts, stream=True)
            else:
                response = model.generate_content(parts, stream=True)
            for chunk in response:
                if chunk.text:
                    q.put(("text", chunk.text))
            q.put(("done", None))
        except Exception as e:
            q.put(("error", str(e)))
    
    # Start generation in background thread
    thread = threading.Thread(target=_sync_stream)
    thread.start()
    
    # Yield chunks as they arrive
    while True:
        try:
            # Use a timeout to allow asyncio to breathe
            item = await loop.run_in_executor(None, lambda: q.get(timeout=0.1))
            event_type, data = item
            if event_type == "done":
                break
            elif event_type == "error":
                yield f"\n\nError: {data}"
                break
            elif event_type == "text":
                yield data
        except:
            # Timeout, continue waiting
            if not thread.is_alive():
                break
            await asyncio.sleep(0.01)
    
    thread.join()
