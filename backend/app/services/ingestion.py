import asyncio
from datetime import datetime
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession, create_async_engine, async_sessionmaker
import uuid
import io
import json
import csv

from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from app.config import get_settings
from app.db.models import Folder, FolderStatus, IndexMode
from app.services.google_drive import GoogleDriveService, SUPPORTED_MIME_TYPES
from app.services import vector_store
from app.services.gemini_service import (
    should_use_fast_path,
    upload_files_to_gemini,
    MAX_FILES_FOR_FAST_PATH,
    MAX_TOTAL_SIZE_BYTES,
)

settings = get_settings()

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap
    
    return chunks


def extract_text_from_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def extract_pdf_pages(content: bytes) -> list[tuple[int, str]]:
    reader = PdfReader(io.BytesIO(content))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((i, text))
    return pages


def extract_text_from_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), data_only=True)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text += row_text + "\n"
    return text


def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts)


def extract_text_from_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        rows.append(" | ".join(row))
    return "\n".join(rows)


def extract_text_from_json(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    try:
        data = json.loads(text)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return text


TEXT_MIME_TYPES = {
    "text/plain", "text/markdown", "text/html", "application/rtf", "text/rtf",
    "text/x-python", "application/x-python-code", "text/javascript", 
    "application/javascript", "text/x-java-source", "text/x-c", "text/x-c++src",
    "text/x-csharp", "text/x-go", "text/x-rust", "text/x-typescript",
    "application/xml", "text/xml", "application/x-yaml", "text/yaml",
}


def extract_text(content: bytes, mime_type: str) -> str:
    if mime_type == "application/pdf":
        return extract_text_from_pdf(content)
    elif mime_type in [
        "application/vnd.google-apps.document",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ]:
        return extract_text_from_docx(content)
    elif mime_type == "application/msword":
        try:
            return extract_text_from_docx(content)
        except Exception:
            return content.decode("utf-8", errors="ignore")
    elif mime_type in [
        "application/vnd.google-apps.spreadsheet",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ]:
        return extract_text_from_xlsx(content)
    elif mime_type == "application/vnd.ms-excel":
        try:
            return extract_text_from_xlsx(content)
        except Exception:
            return ""
    elif mime_type in [
        "application/vnd.google-apps.presentation",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ]:
        return extract_text_from_pptx(content)
    elif mime_type == "application/vnd.ms-powerpoint":
        try:
            return extract_text_from_pptx(content)
        except Exception:
            return ""
    elif mime_type == "text/csv":
        return extract_text_from_csv(content)
    elif mime_type == "application/json":
        return extract_text_from_json(content)
    elif mime_type in TEXT_MIME_TYPES:
        return content.decode("utf-8", errors="ignore")
    return ""


def process_file_content(file: dict, content: bytes) -> tuple[list[str], list[dict], list[str]]:
    mime_type = file["mimeType"]
    
    if mime_type == "application/pdf":
        return process_pdf_content(file, content)
    
    text = extract_text(content, mime_type)
    if not text.strip():
        return [], [], []
    
    chunks = chunk_text(text)
    documents = []
    metadatas = []
    ids = []
    
    for i, chunk in enumerate(chunks):
        documents.append(chunk)
        metadatas.append({
            "file_id": file["id"],
            "file_name": file["name"],
            "file_path": file.get("path", file["name"]),
            "mime_type": mime_type,
            "chunk_index": i,
            "total_chunks": len(chunks),
        })
        ids.append(f"{file['id']}_chunk_{i}")
    
    return documents, metadatas, ids


def process_pdf_content(file: dict, content: bytes) -> tuple[list[str], list[dict], list[str]]:
    pages = extract_pdf_pages(content)
    if not pages:
        return [], [], []
    
    documents = []
    metadatas = []
    ids = []
    chunk_idx = 0
    
    for page_num, page_text in pages:
        page_chunks = chunk_text(page_text)
        for i, chunk in enumerate(page_chunks):
            documents.append(chunk)
            metadatas.append({
                "file_id": file["id"],
                "file_name": file["name"],
                "file_path": file.get("path", file["name"]),
                "mime_type": "application/pdf",
                "chunk_index": chunk_idx,
                "page_number": page_num,
                "total_pages": len(pages),
                "chunk_in_page": i,
            })
            ids.append(f"{file['id']}_page_{page_num}_chunk_{i}")
            chunk_idx += 1
    
    return documents, metadatas, ids


async def process_files_parallel(downloaded: list[tuple[dict, bytes | None]]) -> tuple[list[str], list[dict], list[str], int]:
    loop = asyncio.get_running_loop()
    
    def process_batch(items):
        results = []
        for file, content in items:
            if content is None:
                results.append(([], [], [], False))
                continue
            docs, metas, ids = process_file_content(file, content)
            results.append((docs, metas, ids, bool(docs)))
        return results
    
    from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor
    
    batch_size = max(1, len(downloaded) // 4)
    batches = [downloaded[i:i + batch_size] for i in range(0, len(downloaded), batch_size)]
    
    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = [loop.run_in_executor(executor, process_batch, batch) for batch in batches]
        batch_results = await asyncio.gather(*futures)
    
    all_documents = []
    all_metadatas = []
    all_ids = []
    indexed_count = 0
    
    for batch_result in batch_results:
        for docs, metas, ids, success in batch_result:
            if success:
                all_documents.extend(docs)
                all_metadatas.extend(metas)
                all_ids.extend(ids)
                indexed_count += 1
    
    return all_documents, all_metadatas, all_ids, indexed_count


async def ingest_folder(folder_id: str, user_id: str, refresh_token: str):
    """
    Ingest a folder using the optimal path:
    - Fast Path (Gemini Files): For small folders (<5 files, <20MB total)
    - RAG Path (Chroma): For larger folders
    """
    engine = create_async_engine(
        settings.database_url.replace("postgresql://", "postgresql+asyncpg://"),
        echo=False,
    )
    async_session = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)
    
    async with async_session() as db:
        result = await db.execute(select(Folder).where(Folder.id == uuid.UUID(folder_id)))
        folder = result.scalar_one_or_none()
        if not folder:
            return
        
        folder.status = FolderStatus.INDEXING
        await db.commit()
        
        try:
            print(f"[Ingest] Starting ingestion for drive_folder_id: {folder.drive_folder_id}")
            drive_service = GoogleDriveService(refresh_token)
            files = await drive_service.list_files(folder.drive_folder_id, recursive=True)
            
            file_manifest = [
                {
                    "id": f["id"],
                    "name": f["name"],
                    "path": f.get("path", f["name"]),
                    "mime_type": f["mimeType"],
                    "size": f.get("size", 0),
                }
                for f in files
            ]
            
            if not files:
                print(f"[Ingest] WARNING: No supported files found in folder!")
                folder.status = FolderStatus.READY
                folder.file_count = 0
                folder.indexed_at = datetime.utcnow()
                await db.commit()
                return
            
            use_fast_path = should_use_fast_path(files)
            print(f"[Ingest] Found {len(files)} files. Using {'FAST PATH (Gemini Files)' if use_fast_path else 'RAG PATH (Chroma)'}")
            
            print(f"[Ingest] Starting download of {len(files)} files...")
            downloaded = await drive_service.download_files_parallel(files, max_concurrent=10)
            print(f"[Ingest] Downloaded {len([d for d in downloaded if d[1] is not None])} files successfully")
            
            if use_fast_path:
                try:
                    gemini_files = await upload_files_to_gemini(downloaded)
                    successful_uploads = [f for f in gemini_files if f.get("gemini_uri")]
                    total_attempted = len(gemini_files)
                    success_rate = len(successful_uploads) / total_attempted if total_attempted > 0 else 0
                    
                    if success_rate >= 0.5 and successful_uploads:
                        fast_manifest = [
                            {
                                "id": f.get("id", ""),
                                "name": f.get("name", ""),
                                "path": f.get("name", ""),
                                "mime_type": f.get("mime_type", ""),
                                "size": f.get("size", 0),
                            }
                            for f in successful_uploads
                        ]
                        vector_store.store_file_manifest(user_id, folder_id, fast_manifest)
                        
                        folder.index_mode = IndexMode.gemini_files
                        folder.gemini_files = successful_uploads
                        folder.status = FolderStatus.READY
                        folder.file_count = len(successful_uploads)
                        folder.indexed_at = datetime.utcnow()
                        await db.commit()
                        print(f"[Ingest] FAST PATH complete: {folder.file_count}/{total_attempted} files uploaded to Gemini")
                        return
                    else:
                        print(f"[Ingest] FAST PATH failed - only {len(successful_uploads)}/{total_attempted} uploads succeeded ({success_rate*100:.0f}%), falling back to RAG")
                except Exception as e:
                    print(f"[Ingest] FAST PATH failed with error: {e}, falling back to RAG")
            
            # RAG PATH (also used as fallback if fast path fails)
            all_documents, all_metadatas, all_ids, indexed_count = await process_files_parallel(downloaded)
            print(f"[Ingest] RAG PATH: Indexed {indexed_count} files, {len(all_documents)} chunks")
            
            if all_documents:
                vector_store.add_documents(
                    user_id=user_id,
                    folder_id=folder_id,
                    documents=all_documents,
                    metadatas=all_metadatas,
                    ids=all_ids,
                )
            
            vector_store.store_file_manifest(user_id, folder_id, file_manifest)
            
            folder.index_mode = IndexMode.chroma
            folder.status = FolderStatus.READY
            folder.file_count = indexed_count
            folder.indexed_at = datetime.utcnow()
            await db.commit()
            
        except Exception as e:
            print(f"[Ingest] FAILED: {e}")
            folder.status = FolderStatus.FAILED
            await db.commit()
            raise e


async def ingest_single_file(folder_id: str, user_id: str, refresh_token: str, file_id: str, file_metadata: dict):
    """
    Ingest a single file - always uses FAST PATH (Gemini Files) for speed.
    """
    engine = create_async_engine(
        settings.database_url.replace("postgresql://", "postgresql+asyncpg://"),
        echo=False,
    )
    async_session = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)
    
    async with async_session() as db:
        result = await db.execute(select(Folder).where(Folder.id == uuid.UUID(folder_id)))
        folder = result.scalar_one_or_none()
        if not folder:
            return
        
        folder.status = FolderStatus.INDEXING
        await db.commit()
        
        try:
            drive_service = GoogleDriveService(refresh_token)
            mime_type = file_metadata.get("mimeType", "")
            
            content = await drive_service.download_file(file_id, mime_type)
            file_size = len(content) if content else 0
            
            print(f"[Ingest] Single file: {file_metadata['name']} ({file_size / 1024:.1f} KB)")
            
            file_dict = {
                "id": file_id,
                "name": file_metadata["name"],
                "path": file_metadata["name"],
                "mimeType": mime_type,
                "size": file_size,
            }
            
            if file_size < 10 * 1024 * 1024:
                print(f"[Ingest] Using FAST PATH (Gemini Files)")
                gemini_files = await upload_files_to_gemini([(file_dict, content)])
                
                if gemini_files and gemini_files[0].get("gemini_uri"):
                    single_manifest = [{
                        "id": file_id,
                        "name": file_metadata["name"],
                        "path": file_metadata["name"],
                        "mime_type": mime_type,
                        "size": file_size,
                    }]
                    vector_store.store_file_manifest(user_id, folder_id, single_manifest)
                    
                    folder.index_mode = IndexMode.gemini_files
                    folder.gemini_files = gemini_files
                    folder.status = FolderStatus.READY
                    folder.file_count = 1
                    folder.indexed_at = datetime.utcnow()
                    await db.commit()
                    print(f"[Ingest] FAST PATH complete")
                    return
                else:
                    print(f"[Ingest] Gemini upload failed, falling back to Chroma")
            
            # Fallback to RAG PATH for large files or if Gemini upload failed
            print(f"[Ingest] Using RAG PATH (Chroma)")
            documents, metadatas, ids = process_file_content(file_dict, content)
            
            if not documents:
                folder.status = FolderStatus.FAILED
                await db.commit()
                return
            
            vector_store.add_documents(
                user_id=user_id,
                folder_id=folder_id,
                documents=documents,
                metadatas=metadatas,
                ids=ids,
            )
            
            file_manifest = [{
                "id": file_id,
                "name": file_metadata["name"],
                "path": file_metadata["name"],
                "mime_type": mime_type,
                "size": file_size,
            }]
            vector_store.store_file_manifest(user_id, folder_id, file_manifest)
            
            folder.index_mode = IndexMode.chroma
            folder.status = FolderStatus.READY
            folder.file_count = 1
            folder.indexed_at = datetime.utcnow()
            await db.commit()
            print(f"[Ingest] RAG PATH complete")
            
        except Exception as e:
            print(f"[Ingest] FAILED: {e}")
            folder.status = FolderStatus.FAILED
            await db.commit()
            raise e
